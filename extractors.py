"""
Content extractors for different file types in mtexts.

This module provides functions to extract text content from various file types
including PDFs, Google Docs, PowerPoint files, etc.
"""

import io
import logging
import tempfile
import os
from pathlib import Path

# Third-party imports for handling different file types
import PyPDF2
try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pptx
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BS4_AVAILABLE = False

# Import the Google Drive functions
from google_drive import export_google_doc

# Configure logging
logger = logging.getLogger(__name__)


def extract_file_content(drive_client, file_info):
    """
    Extract content from a file based on its MIME type.
    
    Args:
        drive_client: Google Drive API service object
        file_info: File metadata dictionary
    
    Returns:
        String containing the extracted text content
    """
    mime_type = file_info.get('mimeType', '')
    file_id = file_info.get('id')
    file_name = file_info.get('name', 'Unknown file')
    
    logger.info(f"Extracting content from: {file_name} (MIME: {mime_type})")
    
    # Handle Google Workspace documents
    if mime_type == 'application/vnd.google-apps.document':
        return extract_google_doc(drive_client, file_id)
    elif mime_type == 'application/vnd.google-apps.spreadsheet':
        return extract_google_sheet(drive_client, file_id)
    elif mime_type == 'application/vnd.google-apps.presentation':
        return extract_google_slides(drive_client, file_id)
        
    # Handle PDFs
    elif mime_type == 'application/pdf':
        return extract_pdf(drive_client, file_id)
        
    # Handle Office documents
    elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
        return extract_docx(drive_client, file_id)
    elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        return extract_pptx(drive_client, file_id)
        
    # Handle text-based formats
    elif mime_type == 'text/plain':
        return extract_text_file(drive_client, file_id)
    elif mime_type == 'text/html':
        return extract_html(drive_client, file_id)
    elif mime_type == 'text/markdown':
        return extract_text_file(drive_client, file_id)
        
    # Unsupported file type
    else:
        logger.warning(f"Unsupported MIME type: {mime_type}")
        return f"[Content extraction not supported for file type: {mime_type}]"


def extract_google_doc(drive_client, file_id):
    """Extract text from a Google Doc."""
    try:
        content = export_google_doc(drive_client, file_id, mime_type='text/plain')
        if content:
            return content.getvalue().decode('utf-8')
        return ""
    except Exception as e:
        logger.error(f"Error extracting Google Doc: {e}")
        return f"[Error extracting Google Doc: {str(e)}]"


def extract_google_sheet(drive_client, file_id):
    """Extract text from a Google Sheet."""
    try:
        content = export_google_doc(drive_client, file_id, mime_type='text/csv')
        if content:
            return content.getvalue().decode('utf-8')
        return ""
    except Exception as e:
        logger.error(f"Error extracting Google Sheet: {e}")
        return f"[Error extracting Google Sheet: {str(e)}]"


def extract_google_slides(drive_client, file_id):
    """Extract text from Google Slides."""
    try:
        # Export as plain text first
        content = export_google_doc(drive_client, file_id, mime_type='text/plain')
        if content:
            return content.getvalue().decode('utf-8')
            
        # Fallback: try to export as PDF and extract text
        pdf_content = export_google_doc(drive_client, file_id, mime_type='application/pdf')
        if pdf_content:
            return extract_text_from_pdf_bytes(pdf_content)
            
        return ""
    except Exception as e:
        logger.error(f"Error extracting Google Slides: {e}")
        return f"[Error extracting Google Slides: {str(e)}]"


def extract_pdf(drive_client, file_id):
    """Extract text from a PDF file."""
    try:
        from google_drive import download_file
        content = download_file(drive_client, file_id)
        if content:
            return extract_text_from_pdf_bytes(content)
        return ""
    except Exception as e:
        logger.error(f"Error extracting PDF: {e}")
        return f"[Error extracting PDF: {str(e)}]"


def extract_text_from_pdf_bytes(pdf_bytes):
    """Extract text from PDF file bytes."""
    text = ""
    try:
        pdf_reader = PyPDF2.PdfReader(pdf_bytes)
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text += page.extract_text() + "\n\n"
        return text
    except Exception as e:
        logger.error(f"Error extracting text from PDF bytes: {e}")
        return f"[Error extracting PDF content: {str(e)}]"


def extract_docx(drive_client, file_id):
    """Extract text from a DOCX file."""
    if not DOCX_AVAILABLE:
        return "[python-docx library not installed. Cannot extract DOCX content.]"
        
    try:
        from google_drive import download_file
        content = download_file(drive_client, file_id)
        if not content:
            return ""
            
        # Create a temporary file to load the docx
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_file:
            temp_file.write(content.getvalue())
            temp_path = temp_file.name
            
        # Extract text from the docx
        doc = docx.Document(temp_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
            
        # Clean up temporary file
        try:
            os.unlink(temp_path)
        except:
            pass
            
        return '\n'.join(full_text)
    except Exception as e:
        logger.error(f"Error extracting DOCX: {e}")
        return f"[Error extracting DOCX: {str(e)}]"


def extract_pptx(drive_client, file_id):
    """Extract text from a PPTX file."""
    if not PPTX_AVAILABLE:
        return "[python-pptx library not installed. Cannot extract PPTX content.]"
        
    try:
        from google_drive import download_file
        content = download_file(drive_client, file_id)
        if not content:
            return ""
            
        # Create a temporary file to load the pptx
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
            temp_file.write(content.getvalue())
            temp_path = temp_file.name
            
        # Extract text from the pptx
        presentation = pptx.Presentation(temp_path)
        full_text = []
        
        # Extract text from slides
        for slide in presentation.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            if slide_text:
                full_text.append("\n".join(slide_text))
                
        # Clean up temporary file
        try:
            os.unlink(temp_path)
        except:
            pass
            
        return '\n\n'.join(full_text)
    except Exception as e:
        logger.error(f"Error extracting PPTX: {e}")
        return f"[Error extracting PPTX: {str(e)}]"


def extract_text_file(drive_client, file_id):
    """Extract text from a plain text file."""
    try:
        from google_drive import download_file
        content = download_file(drive_client, file_id)
        if content:
            return content.getvalue().decode('utf-8', errors='replace')
        return ""
    except Exception as e:
        logger.error(f"Error extracting text file: {e}")
        return f"[Error extracting text file: {str(e)}]"


def extract_html(drive_client, file_id):
    """Extract text from an HTML file."""
    if not BS4_AVAILABLE:
        return "[BeautifulSoup4 library not installed. Cannot extract HTML content.]"
        
    try:
        from google_drive import download_file
        content = download_file(drive_client, file_id)
        if not content:
            return ""
            
        html_content = content.getvalue().decode('utf-8', errors='replace')
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.extract()
            
        # Get text
        text = soup.get_text()
        
        # Break into lines and remove leading and trailing space
        lines = (line.strip() for line in text.splitlines())
        # Break multi-headlines into a line each
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        # Drop blank lines
        text = '\n'.join(chunk for chunk in chunks if chunk)
        
        return text
    except Exception as e:
        logger.error(f"Error extracting HTML: {e}")
        return f"[Error extracting HTML: {str(e)}]"